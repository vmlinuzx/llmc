"""
Document Sidecar System - Convert binary documents to markdown for RAG indexing.

This module handles conversion of binary document formats (PDF, DOCX, etc.) to
gzipped markdown sidecars. The sidecars are stored in .llmc/sidecars/ and indexed
using the standard TechDocsExtractor pipeline.

Key features:
- Transparent conversion: PDFs → markdown → gzip
- Staleness detection: re-convert only when source changes
- Path preservation: search results show original PDF path, not sidecar
- Lifecycle management: clean up orphaned sidecars

Usage:
    converter = SidecarConverter()
    sidecar_path = converter.convert(source_path, repo_root)
    
    if converter.is_stale(source_path, repo_root):
        converter.convert(source_path, repo_root)
"""

from __future__ import annotations

import gzip
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    pass

log = logging.getLogger(__name__)

# Supported sidecar extensions
SIDECAR_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx", ".rtf"})


def is_sidecar_eligible(path: Path) -> bool:
    """Check if a file should have a sidecar generated."""
    return path.suffix.lower() in SIDECAR_EXTENSIONS


def get_sidecar_path(source: Path, repo_root: Path) -> Path:
    """Compute sidecar path: .llmc/sidecars/<rel_path>.md.gz"""
    try:
        rel = source.relative_to(repo_root)
    except ValueError:
        # source is already relative
        rel = source
    return repo_root / ".llmc" / "sidecars" / f"{rel}.md.gz"


def is_sidecar_stale(source: Path, repo_root: Path) -> bool:
    """Check if sidecar needs regeneration.
    
    Returns True if:
    - Sidecar doesn't exist
    - Source is newer than sidecar (mtime comparison)
    """
    sidecar = get_sidecar_path(source, repo_root)
    
    if not sidecar.exists():
        return True
    
    try:
        source_abs = repo_root / source if not source.is_absolute() else source
        if source_abs.stat().st_mtime > sidecar.stat().st_mtime:
            return True
    except OSError:
        return True  # If we can't stat, assume stale
    
    return False


def cleanup_orphan_sidecars(repo_root: Path) -> int:
    """Remove sidecars whose source files no longer exist.
    
    Returns:
        Number of orphaned sidecars removed
    """
    sidecars_dir = repo_root / ".llmc" / "sidecars"
    if not sidecars_dir.exists():
        return 0
    
    removed = 0
    for sidecar in sidecars_dir.rglob("*.md.gz"):
        # Reconstruct source path from sidecar path
        try:
            rel_sidecar = sidecar.relative_to(sidecars_dir)
            # Remove .md.gz suffix to get original relative path
            source_rel = str(rel_sidecar).removesuffix(".md.gz")
            source_path = repo_root / source_rel
            
            if not source_path.exists():
                sidecar.unlink()
                removed += 1
                log.info(f"Removed orphan sidecar: {sidecar}")
                
                # Also try to remove empty parent directories
                try:
                    sidecar.parent.rmdir()
                except OSError:
                    pass  # Directory not empty, that's fine
        except Exception as e:
            log.warning(f"Error checking sidecar {sidecar}: {e}")
    
    return removed


class DocumentConverter(ABC):
    """Base class for document-to-markdown converters."""
    
    @abstractmethod
    def convert(self, path: Path) -> str:
        """Convert document to markdown.
        
        Args:
            path: Path to the source document
            
        Returns:
            Markdown content as string
        """
        pass


class PdfToMarkdown(DocumentConverter):
    """Convert PDF to markdown using PyMuPDF (fitz)."""
    
    def convert(self, path: Path) -> str:
        try:
            import fitz  # pymupdf
        except ImportError:
            raise ImportError(
                "pymupdf is required for PDF conversion. "
                "Install with: pip install pymupdf"
            )
        
        doc = fitz.open(path)
        lines = []
        
        # Add document title if available
        metadata = doc.metadata
        if metadata and metadata.get("title"):
            lines.append(f"# {metadata['title']}\n")
        else:
            lines.append(f"# {path.stem}\n")
        
        lines.append(f"\n*Source: {path.name}*\n")
        
        try:
            for page_num, page in enumerate(doc, 1):
                # Add page marker as H2 heading
                lines.append(f"\n## Page {page_num}\n")
                
                # Extract text blocks (preserves reading order better than raw text)
                blocks = page.get_text("blocks")
                for block in blocks:
                    # block format: (x0, y0, x1, y1, text, block_no, block_type)
                    if len(block) >= 5:
                        text = str(block[4]).strip()
                        if text:
                            # Clean up common PDF artifacts
                            text = text.replace("\x00", "")  # Null bytes
                            lines.append(text + "\n")
        finally:
            doc.close()
        
        return "\n".join(lines)


class DocxToMarkdown(DocumentConverter):
    """Convert DOCX to markdown preserving structure."""
    
    def convert(self, path: Path) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX conversion. "
                "Install with: pip install python-docx"
            )
        
        doc = Document(path)
        lines = []
        
        # Add document title
        lines.append(f"# {path.stem}\n")
        lines.append(f"\n*Source: {path.name}*\n")
        
        for para in doc.paragraphs:
            style_name = (para.style.name or "").lower()
            text = para.text.strip()
            
            if not text:
                continue
            
            # Convert headings based on style
            if "heading 1" in style_name:
                lines.append(f"\n## {text}\n")
            elif "heading 2" in style_name:
                lines.append(f"\n### {text}\n")
            elif "heading 3" in style_name:
                lines.append(f"\n#### {text}\n")
            elif "title" in style_name:
                lines.append(f"\n# {text}\n")
            else:
                lines.append(f"{text}\n")
        
        return "\n".join(lines)


class PptxToMarkdown(DocumentConverter):
    """Convert PPTX to markdown with slide-per-section structure."""
    
    def convert(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX conversion. "
                "Install with: pip install python-pptx"
            )
        
        prs = Presentation(path)
        lines = []
        
        # Add document title
        lines.append(f"# {path.stem}\n")
        lines.append(f"\n*Source: {path.name}*\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            lines.append(f"\n## Slide {slide_num}\n")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        lines.append(f"{text}\n")
            
            # Include speaker notes if present
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"\n*Speaker notes: {notes}*\n")
        
        return "\n".join(lines)


class RtfToMarkdown(DocumentConverter):
    """Convert RTF to markdown (basic extraction)."""
    
    def convert(self, path: Path) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise ImportError(
                "striprtf is required for RTF conversion. "
                "Install with: pip install striprtf"
            )
        
        content = path.read_text(encoding="utf-8", errors="ignore")
        text = rtf_to_text(content)
        
        lines = []
        lines.append(f"# {path.stem}\n")
        lines.append(f"\n*Source: {path.name}*\n\n")
        lines.append(text)
        
        return "\n".join(lines)


class SidecarConverter:
    """Convert binary documents to gzipped markdown sidecars."""
    
    CONVERTERS: dict[str, type[DocumentConverter]] = {
        ".pdf": PdfToMarkdown,
        ".docx": DocxToMarkdown,
        ".pptx": PptxToMarkdown,
        ".rtf": RtfToMarkdown,
    }
    
    def __init__(self):
        self._converter_cache: dict[str, DocumentConverter] = {}
    
    def _get_converter(self, ext: str) -> DocumentConverter | None:
        """Get or create a converter for the given extension."""
        ext = ext.lower()
        if ext not in self.CONVERTERS:
            return None
        
        if ext not in self._converter_cache:
            self._converter_cache[ext] = self.CONVERTERS[ext]()
        
        return self._converter_cache[ext]
    
    def can_convert(self, path: Path) -> bool:
        """Check if this file can be converted to a sidecar."""
        return path.suffix.lower() in self.CONVERTERS
    
    def convert(self, source_path: Path, repo_root: Path) -> Path | None:
        """Convert source document to gzipped markdown sidecar.
        
        Args:
            source_path: Path to source document (relative or absolute)
            repo_root: Repository root path
            
        Returns:
            Path to generated sidecar, or None if conversion failed/unsupported
        """
        # Ensure we have an absolute source path
        if not source_path.is_absolute():
            source_abs = repo_root / source_path
        else:
            source_abs = source_path
            source_path = source_path.relative_to(repo_root)
        
        if not source_abs.exists():
            log.warning(f"Source file not found: {source_abs}")
            return None
        
        ext = source_abs.suffix.lower()
        converter = self._get_converter(ext)
        
        if converter is None:
            log.debug(f"No converter for extension: {ext}")
            return None
        
        try:
            markdown = converter.convert(source_abs)
        except ImportError as e:
            log.warning(f"Missing dependency for {ext} conversion: {e}")
            return None
        except Exception as e:
            log.error(f"Failed to convert {source_path}: {e}")
            return None
        
        sidecar_path = get_sidecar_path(source_path, repo_root)
        sidecar_path.parent.mkdir(parents=True, exist_ok=True)
        
        try:
            with gzip.open(sidecar_path, "wt", encoding="utf-8") as f:
                f.write(markdown)
            log.info(f"Generated sidecar: {sidecar_path}")
            return sidecar_path
        except Exception as e:
            log.error(f"Failed to write sidecar {sidecar_path}: {e}")
            return None
    
    def read_sidecar(self, sidecar_path: Path) -> str | None:
        """Read content from a gzipped sidecar file.
        
        Args:
            sidecar_path: Path to the .md.gz sidecar
            
        Returns:
            Markdown content as string, or None if read failed
        """
        try:
            with gzip.open(sidecar_path, "rt", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            log.error(f"Failed to read sidecar {sidecar_path}: {e}")
            return None


# Singleton for convenience
_default_converter: SidecarConverter | None = None


def get_converter() -> SidecarConverter:
    """Get the default sidecar converter instance."""
    global _default_converter
    if _default_converter is None:
        _default_converter = SidecarConverter()
    return _default_converter
